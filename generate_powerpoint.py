import tictactoe
from pptx import Presentation
from pptx.util import Inches, Pt

states, state_graph = tictactoe.compute_graph()

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

player_markers = {
    tictactoe.PLAYER_X: 'X',
    tictactoe.PLAYER_O: 'O',
    tictactoe.BLANK: '_'
}

state_shapes = {}

# Create the slides
for state in states:
    slide = prs.slides.add_slide(blank_slide_layout)
    player, board = state

    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tf = box.text_frame
    tf.text = "%s’s turn" % player_markers[player]

    row_shapes = []
    for i in range(len(board)):
        col_shapes = []
        for j in range(len(board[i])):
            box = slide.shapes.add_textbox(Inches(2.5 + j), Inches(1.75 + i), Inches(5/3), Inches(5/3))
            tf = box.text_frame
            tf.text = player_markers[board[i][j]]
            col_shapes.append(box)
        row_shapes.append(col_shapes)
    state_shapes[state] = row_shapes

# Add the links between squares and slides for interactivity.
for (state, neighbors) in state_graph.items():
    for (i, j), index in neighbors:
        state_shapes[state][i][j].click_action.target_slide = prs.slides[index]

prs.save('tictactoe.pptx')
